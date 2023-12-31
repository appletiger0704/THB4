# -*- coding: utf-8 -*-
"""
Created on Sun Dec 17 14:40:48 2023

@author: User
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import pandas as pd 
import re
import os

# 日期物件
now = datetime.now()
yesterday = now - timedelta(days = 1)
next_day = now + timedelta(days = 1)
day_after_tomorrow = now + timedelta(days = 2)

# 日期
format_date = now.strftime("%Y%m%d")
format_ydate = yesterday.strftime("%Y%m%d")

# 日
yesterday_day = yesterday.strftime("%d")
now_day = now.strftime("%d")
next_day = next_day.strftime("%d")
after_tomorrow_day = day_after_tomorrow.strftime("%d")

# 月
now_month = now.strftime("%m")
yesterday_month = yesterday.strftime("%m")
after_tomorrow_month = day_after_tomorrow.strftime("%m")


# 程式執行位置
path = rf"C:\Users\User\Desktop\East_auto\{format_date}"
# 今日檔案位置
ppt_path = rf"C:\Users\User\Desktop\East_auto\{format_date}\東分局-{format_date}-0900-AM-常時分析報告.pptx"
# 昨日檔案位置
ppt_ypath = rf"C:\Users\User\Desktop\East_auto\{format_ydate}\東分局-{format_ydate}-0900-AM-常時分析報告.pptx"
# 昨日雨量csv檔案位置
csv_path = rf"C:\Users\User\Desktop\East_auto\yday_accumulate\{format_date}\{format_ydate}_累積雨量.csv"
os.chdir(path)



def convert_to_minguo(year):
    
    return year - 1911



def write_txt(text):
    date = now.strftime("%y%m%d")
    with open(os.path.join(path, f"{date}_img.txt"), 'a') as file:
        file.write(text + "\n")


# 字型設定
def font(textbox, Date_String, RGB, size):
    
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.clear()
    run = paragraph.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = RGB
    run.text = Date_String
    run.font.bold = True
    run.font.name = "微軟正黑體"
    paragraph.alignment = PP_ALIGN.CENTER


# 更新文字方塊日期
def update_date(slide, re_express, Date_String, RGB, size):
    pattern = re.compile(re_express)
    
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                raw_text = shape.text_frame.text
                # is "raw_text" match pattern
                if re.search(pattern, raw_text):
                    font(shape, Date_String, RGB, size)
                    break 
                
        except Exception as e:
            write_txt(f"{slide} update have unexpected exception: {e}")


# 更新表格內日期
def table_update_date(slide, start, end, RGB, size):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for i in range(start, end):
                slide_text_box = shape.table.cell(0, i)
                if i == start:
                    Date_String = f"今({now_day})日"
                elif i == (start + 1):
                    Date_String = f"明({next_day})日"
                elif i == (start + 2):
                    Date_String = f"後({after_tomorrow_day})日"
                font(slide_text_box, Date_String, RGB, size)
                
                
def road_table_update_date(slide, start, end, RGB, size):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for i in range(start, end):
                slide_text_box = shape.table.cell(0, i)
                if i == start:
                    Date_String = f"昨({yesterday_day})日雨量累積"
                elif i == (start + 1):
                    Date_String = f"今({now_day})日"
                elif i == (start + 2):
                    Date_String = f"明({next_day})日"
                elif i == (start + 3):
                    Date_String = f"後({after_tomorrow_day})日"
                slide_text_box.text = ""
                font(slide_text_box, Date_String, RGB, size)            


try:
    if os.path.exists(ppt_path):
        
        write_txt("採用今日簡報")
        prs = Presentation(ppt_path)
    
    else:
        
        write_txt("採用昨日簡報")
        prs = Presentation(ppt_ypath)
    
except Exception as e:
    write_txt(f"exception occur when open ppt, error：{e}")
    


first_page_pattern = r"發佈時間：\d{3}年\d{1,2}月\d{1,2}日 \d{4}時"
SWM_pattern = r"\d{1,2}月\d{1,2}日 02:00 地面天氣圖"
Satellite_pattern = r"\d{1,2}月\d{1,2}日 \d{1,2}:\d{1,2} 衛星雲圖"
# StreamLine_pattern = r"\d{1,2}月\d{1,2}日 05:00 700-850hPa 平均駛流場圖"
yday_accumulate_pattern = r"昨\(\d{1,2}\)日 累積雨量"
tday_accumulate_pattern = r"今\(\d{1,2}\)日 00-\d{2}時 累積雨量"
today_QPF_pattern = r"今\(\d{1,2}\)日$"
tomorrow_QPF_pattern = r"明\(\d{1,2}\)日$"
day_after_tomorrow_QPF_pattern = r"後\(\d{1,2}\)日$"


first_slide = prs.slides[0]
second_slide = prs.slides[1]
thrid_slide = prs.slides[2]
forth_slide = prs.slides[3]
fifth_slide = prs.slides[4]
sixth_slide = prs.slides[5]
seventh_slide = prs.slides[6]
eighth_slide = prs.slides[7]
nighth_slide = prs.slides[8]
tenth_slide = prs.slides[9]


update_date(first_slide, first_page_pattern, f"發佈時間：{convert_to_minguo(now.year)}年{now_month}月{now_day}日 0900時", RGBColor(00,00,00), 24)
update_date(second_slide, SWM_pattern, f"{now_month}月{now_day}日 02:00 地面天氣圖", RGBColor(12,51,115), 18)
update_date(second_slide, Satellite_pattern, f"{now_month}月{now_day}日 07:00 衛星雲圖", RGBColor(12,51,115), 18)
# update_date(thrid_slide, StreamLine_pattern, f"{now_month}月{now_day}日 05:00 700-850hPa 平均駛流場圖", RGBColor(12,51,115), 18)
update_date(thrid_slide, yday_accumulate_pattern, f"昨({yesterday_day})日 累積雨量", RGBColor(12,51,115), 18)
update_date(thrid_slide, tday_accumulate_pattern, f"今({now_day})日 00-06時 累積雨量", RGBColor(12,51,115), 18)
update_date(fifth_slide, today_QPF_pattern, f"今({now_day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, tomorrow_QPF_pattern, f"明({next_day})日", RGBColor(12,51,115), 18)
update_date(fifth_slide, day_after_tomorrow_QPF_pattern, f"後({after_tomorrow_day})日", RGBColor(12,51,115), 18)
table_update_date(fifth_slide, 1, 4, RGBColor(00,00,00), 18)
road_table_update_date(seventh_slide, 3, 7, RGBColor(00,00,00), 16)
road_table_update_date(eighth_slide, 3, 7, RGBColor(00,00,00), 16)
table_update_date(nighth_slide, 3, 6, RGBColor(00,00,00), 16)
table_update_date(tenth_slide, 3, 6, RGBColor(00,00,00), 16)


# 更換圖資
def change_img(slide, img, left):
    
    for shape in slide.shapes:
        if shape.shape_type == 13 and (left - 50000<= shape.left <= left + 50000): # 13代表圖片
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # 刪除圖片
            sp = shape._element
            sp.getparent().remove(sp)

            slide.shapes.add_picture(img,  left, top, width, height)
            
            write_txt(f"{img} success  ")
            write_txt(f"{img}, {left}, {top}, {width}, {height} \n")
            print(f"{img} success")
            break



change_img(second_slide, "round_Satellite.png", 6988336)
change_img(second_slide, "round_SWM.png", 1106797)

if os.path.exists(os.path.join(path, "E_06_image.png")):
    
    change_img(thrid_slide, "E_06_image.png", 7066626)
    
else:
    
    change_img(thrid_slide, "E_image.png", 7066626)
    
change_img(thrid_slide, "E_yday_image.png", 1070948)

change_img(forth_slide, "cropped_6_06QPF.png", 1220392)
change_img(forth_slide, "cropped_6_12QPF.png", 3661828)
change_img(forth_slide, "cropped_6_18QPF.png", 6100074)
change_img(forth_slide, "cropped_6_24QPF.png", 8534802)


for i in forth_slide.shapes:
    if i.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print(i.left)


#====================

# 回傳一個list
def read_csv(path):
    
    df = pd.read_csv(path, encoding="big5")

    df.set_index("Unnamed: 0", inplace=True)

    df.index.name = "測站"
    
    df_dict = df.to_dict(orient="dict")["昨日累積雨量"]
    
    rainfall_list = []
    
    for key, item in df_dict.items():
        
        rainfall_list.append(str(item))
    
    return rainfall_list



# 昨日累積雨量
def table_yday_rainfall(slide, start, end):
        
    for shape in slide.shapes:
        
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            
            count = 1
            
            for item in read_csv(csv_path)[start:end]:
                
                rainfall_box = shape.table.cell(count,3)
                font(rainfall_box, item, RGBColor(00,00,00), 16)
                count = count + 1
                

try:

    table_yday_rainfall(seventh_slide, 0, 8)
    table_yday_rainfall(eighth_slide,8, 15)
    print("成功更新監控路段雨量")

except:
    print("未更新監控路段雨量")
    write_txt("未更新監控路段昨日雨量")
    


try:

    prs.save(f'東分局-{format_date}-0900-AM-常時分析報告.pptx')
    write_txt("簡報儲存成功")
    print("成功儲存簡報")

except Exception as e:
    
    write_txt(f"簡報儲存失敗：{e}")
    print("簡報儲存失敗")